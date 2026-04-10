import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(value):
    return RGBColor(*value)


def add_shape(
    slide,
    shape_type,
    left,
    top,
    width,
    height,
    text="",
    fill=(255, 255, 255),
    line=(0, 0, 0),
    font_name="Microsoft YaHei",
    font_size=14,
    bold=False,
    italic=False,
    align=PP_ALIGN.CENTER,
    valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    line_width=1.4,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_name="Microsoft YaHei",
    font_size=12,
    bold=False,
    italic=False,
    align=PP_ALIGN.CENTER,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_connector(slide, x1, y1, x2, y2, dashed=False, width=1.6):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb((0, 0, 0))
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_band(slide, left, top, width, height, text):
    band = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
        text=text,
        fill=(214, 231, 247),
        line=(180, 200, 220),
        font_size=18,
        bold=True,
        line_width=1.0,
    )
    return band


def build_flowchart(output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background panels
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        1.5,
        0.7,
        3.7,
        4.1,
        fill=(245, 250, 255),
        line=(220, 230, 240),
        line_width=1.0,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        5.3,
        0.7,
        3.1,
        4.1,
        fill=(245, 250, 255),
        line=(220, 230, 240),
        line_width=1.0,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        8.5,
        0.7,
        4.6,
        4.1,
        fill=(245, 250, 255),
        line=(220, 230, 240),
        line_width=1.0,
    )

    add_band(slide, 1.5, 0.7, 3.7, 0.35, "1. 输入表示")
    add_band(slide, 5.3, 0.7, 3.1, 0.35, "2. 主干编码")
    add_band(slide, 8.5, 0.7, 4.6, 0.35, "3. 预测与训练")

    # Input block
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        0.15,
        2.2,
        1.55,
        1.7,
        text="输入\n\n物品序列 I\n\n时间序列 T\n\n序列长度 l",
        fill=(250, 253, 255),
        line=(0, 0, 0),
        font_size=15,
        bold=True,
        line_width=1.2,
    )

    # Input representation
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        1.75,
        1.95,
        2.55,
        2.35,
        text="时间感知输入表示",
        fill=(226, 239, 251),
        line=(140, 170, 200),
        font_size=17,
        bold=True,
        line_width=1.2,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        2.05,
        2.55,
        1.9,
        0.72,
        text="CTE：连续/离散时间编码",
        fill=(255, 255, 255),
        line=(110, 110, 110),
        font_size=14,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        2.05,
        3.45,
        1.9,
        0.72,
        text="TGF：时间门控融合",
        fill=(255, 255, 255),
        line=(110, 110, 110),
        font_size=14,
        bold=True,
    )

    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        4.62,
        2.8,
        0.5,
        0.6,
        text="H^(0)",
        fill=(240, 248, 255),
        line=(0, 0, 0),
        font_size=16,
        bold=True,
    )

    # Backbone
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        5.6,
        2.0,
        2.5,
        2.25,
        text="TimeWeaver 主干",
        fill=(226, 239, 251),
        line=(140, 170, 200),
        font_size=17,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        5.95,
        2.65,
        0.9,
        0.85,
        text="上下文流",
        fill=(255, 255, 255),
        line=(110, 110, 110),
        font_size=15,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        5.95,
        3.45,
        0.9,
        0.85,
        text="动态流",
        fill=(255, 255, 255),
        line=(110, 110, 110),
        font_size=15,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        7.15,
        3.0,
        0.75,
        0.65,
        text="融合",
        fill=(210, 244, 244),
        line=(90, 130, 130),
        font_size=15,
        bold=True,
    )

    # Prediction chain
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        8.55,
        3.0,
        1.0,
        0.7,
        text="末位置\n表示",
        fill=(250, 253, 255),
        line=(0, 0, 0),
        font_size=15,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        9.9,
        3.0,
        0.75,
        0.7,
        text="物品\n得分",
        fill=(250, 253, 255),
        line=(0, 0, 0),
        font_size=15,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        10.95,
        3.0,
        0.95,
        0.7,
        text="下一物品\n预测",
        fill=(250, 253, 255),
        line=(0, 0, 0),
        font_size=15,
        bold=True,
    )

    # STC block
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        5.65,
        1.25,
        6.2,
        1.15,
        text="STC（仅训练阶段启用）",
        fill=(247, 247, 240),
        line=(160, 160, 150),
        font_size=17,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        11.0,
        1.28,
        0.72,
        0.28,
        text="仅训练",
        fill=(230, 230, 220),
        line=(130, 130, 120),
        font_size=11,
        bold=False,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        5.9,
        1.62,
        1.05,
        0.55,
        text="调度式时间增强",
        fill=(218, 248, 248),
        line=(90, 130, 130),
        font_size=14,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        7.2,
        1.62,
        0.45,
        0.55,
        text="T_aug",
        fill=(218, 248, 248),
        line=(90, 130, 130),
        font_size=15,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        7.95,
        1.62,
        0.9,
        0.55,
        text="共享\n编码器",
        fill=(218, 248, 248),
        line=(90, 130, 130),
        font_size=14,
        bold=True,
    )
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        9.15,
        1.62,
        1.0,
        0.55,
        text="对比损失",
        fill=(250, 253, 255),
        line=(110, 110, 110),
        font_size=14,
        bold=True,
    )

    # Connectors
    add_connector(slide, 1.70, 3.05, 1.75, 3.05)
    add_connector(slide, 4.30, 3.12, 4.62, 3.12)
    add_connector(slide, 5.12, 3.12, 5.60, 3.12)
    add_connector(slide, 6.85, 3.08, 7.15, 3.08)
    add_connector(slide, 7.90, 3.30, 8.55, 3.30)
    add_connector(slide, 9.55, 3.35, 9.90, 3.35)
    add_connector(slide, 10.65, 3.35, 10.95, 3.35)

    # Inner backbone flow
    add_connector(slide, 6.85, 3.08, 7.15, 3.32)

    # Training dashed flow
    add_connector(slide, 1.72, 2.2, 5.9, 2.2, dashed=True)
    add_connector(slide, 6.95, 1.9, 7.2, 1.9, dashed=True)
    add_connector(slide, 7.65, 1.9, 7.95, 1.9, dashed=True)
    add_connector(slide, 8.85, 1.9, 9.15, 1.9, dashed=True)
    add_connector(slide, 10.15, 2.2, 10.15, 2.8, dashed=True)
    add_connector(slide, 10.15, 2.8, 8.85, 2.8, dashed=True)
    add_connector(slide, 8.85, 2.8, 8.85, 3.0, dashed=True)

    # The line from STC to output branch
    add_connector(slide, 9.85, 2.77, 9.85, 3.0, dashed=True)

    # Labels
    add_textbox(slide, 0.20, 1.85, 1.0, 0.25, "输入", font_size=12, bold=False)
    add_textbox(slide, 11.25, 1.35, 0.85, 0.20, "Training only", font_size=10, bold=False)
    add_textbox(slide, 11.70, 0.95, 0.40, 0.20, "仅训练", font_size=10, bold=False)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate the time-aware flowchart PPT.")
    parser.add_argument(
        "--output-dir",
        default="figures/timeaware_flowchart_ppt",
        help="Directory to save the generated pptx file.",
    )
    parser.add_argument(
        "--filename",
        default="TimeAware_Flowchart_CN.pptx",
        help="Output pptx filename.",
    )
    args = parser.parse_args()

    out_dir = Path(args.output_dir)
    out_path = out_dir / args.filename
    final_path = build_flowchart(out_path)
    print(final_path)


if __name__ == "__main__":
    main()
